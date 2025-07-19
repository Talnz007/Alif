from fastapi import FastAPI, UploadFile, File, HTTPException, Body, APIRouter
from fastapi.responses import JSONResponse
from fastapi.middleware.cors import CORSMiddleware
from typing import List, Dict, Any
import fitz  # PyMuPDF
from io import BytesIO
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain.chains import ConversationalRetrievalChain
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.memory import ConversationBufferMemory
from langchain_community.vectorstores import FAISS
from langchain.docstore.document import Document
from langchain.prompts import PromptTemplate
import os
import uuid
import time
from dotenv import load_dotenv
from pptx import Presentation  # For PPTX support
import requests
from app.core.shared_state import set_vector_db, add_file_metadata, clear_all, get_vector_db, get_uploaded_metadata

router = APIRouter(prefix="/chatbot", tags=["Chatbot"])

# Configuration
USE_VECTOR_SEARCH = os.getenv("USE_VECTOR_SEARCH", "True") == "True"
YOUTUBE_API_KEY = os.getenv("YOUTUBE_API_KEY")
YOUTUBE_API_URL = "https://www.googleapis.com/youtube/v3/search"

chat_history = []

# Load environment variables
load_dotenv()

# Ensure API keys are available
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
if not GOOGLE_API_KEY:
    raise ValueError("GOOGLE_API_KEY environment variable is not set.")

# Initialize FastAPI app
app = FastAPI(title="Alif API", description="AI-powered study companion API", version="1.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# Text extraction functions
def extract_text_from_file(file: UploadFile):
    """Extract text from uploaded file (PDF, TXT, PPTX)."""
    content = file.file.read()
    filename = file.filename

    if file.content_type == "application/pdf":
        doc = fitz.open(stream=content, filetype="pdf")
        return [{"text": page.get_text("text"), "page_number": i + 1} for i, page in enumerate(doc)]
    elif file.content_type == "text/plain":
        return [{"text": content.decode("utf-8"), "page_number": 1}]
    elif file.content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        prs = Presentation(BytesIO(content))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return [{"text": text, "page_number": 1}]
    else:
        raise HTTPException(status_code=400, detail="Unsupported file type")


def process_files(files: List[UploadFile]):
    print(f"process_files called with {len(files)} files: {[f.filename for f in files]}")
    if not files or len(files) == 0:
        print("No files received in process_files.")
        raise HTTPException(status_code=400, detail="No files received")

    clear_all()  # Clear existing vector DB and metadata
    uploaded_files = []
    for file in files:
        print(f"Processing file: {file.filename}, type: {file.content_type}")
        file_size = 0
        try:
            file.file.seek(0, 2)
            file_size = file.file.tell()
            file.file.seek(0)
        except Exception as e:
            print(f"Error getting file size: {e}")

        pages = extract_text_from_file(file)

        file_metadata = {
            "id": str(uuid.uuid4()),
            "filename": file.filename,
            "file_size": file_size,
            "page_count": len(pages),
            "content_type": file.content_type,
            "active": True,
            "created_at": int(time.time())
        }
        add_file_metadata(file_metadata)
        print(f"Added file metadata: {file.filename}, {len(pages)} pages")
        uploaded_files.append(file_metadata)

        if USE_VECTOR_SEARCH:
            all_docs = []
            text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)

            for page in pages:
                chunks = text_splitter.split_text(page["text"])
                for i, chunk in enumerate(chunks):
                    metadata = {
                        "source": file.filename,
                        "page_number": page["page_number"],
                        "chunk_index": i,
                        "file_size": file_size
                    }
                    all_docs.append(Document(page_content=chunk, metadata=metadata))

            if all_docs:
                try:
                    embeddings = GoogleGenerativeAIEmbeddings(
                        model="models/text-embedding-004",
                        google_api_key=GOOGLE_API_KEY
                    )
                    db = FAISS.from_documents(all_docs, embeddings)
                    set_vector_db(db)
                    print(f"Successfully created vector DB with {len(all_docs)} documents")
                except Exception as e:
                    print(f"Warning: Failed to create vector DB but continuing with file metadata: {e}")

    return {
        "success": True,
        "message": "Files processed successfully",
        "data": {"file_count": len(files), "uploaded_files": uploaded_files}
    }


def initialize_qa_chain():
    """Initialize the ConversationalRetrievalChain."""
    vector_db = get_vector_db()
    if vector_db is None:
        raise HTTPException(status_code=400, detail="No documents uploaded yet. Please upload a file first.")

    template = """You are an expert research paper analyst. Use the following pieces of context to provide a detailed answer to the question. If you can't answer based on the context, say so.

Context:
{context}

Question: {question}

Previous conversation:
{chat_history}

Please provide a comprehensive response that:
1. Directly answers the question using information from the uploaded documents
2. Cites specific sources (e.g., document name and page number) when relevant
3. Explains technical concepts clearly
4. Uses examples when helpful
5. Maintains academic accuracy while being accessible

Answer:"""

    prompt = PromptTemplate(
        input_variables=["context", "question", "chat_history"],
        template=template
    )

    memory = ConversationBufferMemory(
        memory_key="chat_history",
        input_key="question",
        output_key="answer",
        return_messages=True
    )

    llm = ChatGoogleGenerativeAI(
        model="gemini-2.0-flash",
        temperature=0.7,
        google_api_key=GOOGLE_API_KEY
    )

    return ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=vector_db.as_retriever(search_kwargs={"k": 3}),
        memory=memory,
        combine_docs_chain_kwargs={"prompt": prompt},
        return_source_documents=True
    )


async def sync_document_context():
    """Sync the document context after processing files."""
    try:
        import httpx
        async with httpx.AsyncClient() as client:
            await client.post(f"http://localhost:8000/api/v1/document-context/sync-from-chatbot")
    except Exception as e:
        print(f"Warning: Could not sync document context: {e}")


async def fetch_youtube_recommendations(query: str, max_results: int = 3):
    """Fetch YouTube video recommendations using the YouTube Data API."""
    if not YOUTUBE_API_KEY:
        raise HTTPException(status_code=500, detail="YouTube API key not configured")

    try:
        params = {
            "part": "snippet",
            "q": query,
            "type": "video",
            "maxResults": max_results,
            "key": YOUTUBE_API_KEY,
        }
        response = requests.get(YOUTUBE_API_URL, params=params)
        response.raise_for_status()
        data = response.json()

        videos = [
            {
                "title": item["snippet"]["title"],
                "url": f"https://www.youtube.com/watch?v={item['id']['videoId']}",
            }
            for item in data.get("items", [])
        ]
        return videos
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to fetch YouTube videos: {e}")


# Endpoints
@router.post("/upload-files/")
async def upload_files(files: List[UploadFile] = File(...)):
    global chat_history
    if not files or len(files) == 0:
        print("No files uploaded in request.")
        return JSONResponse(content={"success": False, "message": "No files uploaded", "data": {"file_count": 0}},
                            status_code=400)
    try:
        result = process_files(files)
        chat_history = []
        await sync_document_context()
        return JSONResponse(content=result)
    except Exception as e:
        import traceback
        error_detail = f"Error processing files: {str(e)}\n{traceback.format_exc()}"
        print(error_detail)
        raise HTTPException(status_code=500, detail=error_detail)


@router.post("/ask-question/")
async def ask_question(data: Dict[str, Any] = Body(...)):
    """Ask a question about the uploaded documents."""
    global chat_history
    question = data.get("question")
    file_ids = data.get("file_ids", [])

    if not question:
        raise HTTPException(status_code=400, detail="Question is required")

    if not get_uploaded_metadata() and USE_VECTOR_SEARCH:
        raise HTTPException(status_code=400, detail="No documents uploaded yet. Please upload a file first.")

    try:
        if USE_VECTOR_SEARCH:
            vector_db = get_vector_db()
            if not vector_db:
                raise HTTPException(status_code=400, detail="No documents uploaded yet. Please upload a file first.")

            qa_chain = initialize_qa_chain()
            response = qa_chain({"question": question})

            answer = response["answer"]
            sources = []
            if "source_documents" in response:
                for doc in response["source_documents"]:
                    source_info = {
                        "title": f"{doc.metadata['source']} (Page {doc.metadata['page_number']})",
                        "content": doc.page_content,
                        "relevance": 0.9
                    }
                    sources.append(source_info)
        else:
            llm = ChatGoogleGenerativeAI(
                model="gemini-2.0-flash",
                temperature=0.7,
                google_api_key=GOOGLE_API_KEY
            )

            uploaded_files = get_uploaded_metadata()
            if not uploaded_files:
                raise HTTPException(status_code=400, detail="No documents uploaded yet. Please upload a file first.")

            file_info = ""
            if uploaded_files:
                file_info = "The user has uploaded the following files:\n" + "\n".join([
                    f"- {f['filename']} ({f['page_count']} pages)" for f in uploaded_files
                    if f["id"] in file_ids
                ])

            conversation_context = ""
            if chat_history:
                conversation_context = "Previous conversation:\n"
                for msg in chat_history[-6:]:
                    role = "User" if msg["role"] == "user" else "Assistant"
                    conversation_context += f"{role}: {msg['content']}\n"

            prompt = f"""You are an expert research paper analyst. Answer the following question:

{question}

{file_info}

{conversation_context}

Note: Advanced document search is currently disabled. Please provide your best answer
based on the information available and general knowledge. If the question seems to require
specific information from the uploaded documents, explain that detailed document search
is temporarily unavailable.

Please provide a comprehensive response that:
1. Answers the question as best as possible with available information
2. Explains technical concepts clearly
3. Uses examples when helpful
4. Maintains academic accuracy while being accessible
"""

            response = llm.invoke(prompt)
            answer = response.content
            sources = []

        chat_history.append({"role": "user", "content": question})
        chat_history.append({"role": "assistant", "content": answer})

        return JSONResponse(content={
            "success": True,
            "message": "Question answered",
            "data": {
                "answer": answer,
                "sources": sources,
                "recommendations": bool(file_ids)  # Indicate if recommendations are available
            }
        })
    except Exception as e:
        import traceback
        error_detail = f"Error processing question: {str(e)}\n{traceback.format_exc()}"
        print(error_detail)
        raise HTTPException(status_code=500, detail=error_detail)


@router.get("/recommendations/")
async def get_recommendations(query: str):
    """Fetch YouTube video recommendations based on query."""
    try:
        videos = await fetch_youtube_recommendations(query)
        return JSONResponse(content={
            "success": True,
            "message": "Recommendations generated",
            "data": {"videos": videos}
        })
    except Exception as e:
        print(f"Error fetching recommendations: {str(e)}")
        return JSONResponse(content={
            "success": False,
            "message": f"Error fetching recommendations: {str(e)}",
            "data": {"videos": []}
        })


@router.get("/chat-history/")
async def get_chat_history():
    """Retrieve the current chat history."""
    return JSONResponse(content={
        "success": True,
        "message": "Chat history retrieved",
        "data": {"chat_history": chat_history}
    })


@router.post("/reset/")
async def reset():
    """Reset the vector store and chat history."""
    global chat_history
    clear_all()
    chat_history = []
    await sync_document_context()
    return JSONResponse(content={
        "success": True,
        "message": "Chatbot reset successfully",
        "data": {}
    })


@router.get("/health/")
async def health_check():
    """Check if the API is running and API key is configured."""
    return JSONResponse(content={
        "success": True,
        "message": "Health check passed",
        "data": {
            "status": "healthy",
            "google_api_key_configured": bool(GOOGLE_API_KEY),
            "youtube_api_key_configured": bool(YOUTUBE_API_KEY),
            "vector_search_enabled": USE_VECTOR_SEARCH
        }
    })